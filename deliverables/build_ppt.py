from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "deliverables" / "A5-灵山小向导-产品方案介绍.pptx"
ASSETS = ROOT / "deliverables" / "assets"
AVATAR = ROOT / "xmov-digital-human-web-demos" / "public" / "assets" / "xmov-male-runtime-cutout.png"
FONT = "Microsoft YaHei"

C = {
    "dark": "1A2E20",
    "section": "233B2A",
    "cream": "F5EDD6",
    "paper": "FFFDF8",
    "white": "FFFFFF",
    "coral": "E8734A",
    "sage": "5B9E8C",
    "teal": "3A8C7A",
    "green": "6AB870",
    "gold": "B78635",
    "ink": "1A2E20",
    "light": "F5EDD6",
    "muted": "728178",
    "divider": "2E4E38",
    "border": "DED3B9",
    "soft": "ECE4D2",
    "red": "A54537",
}


def color(value):
    return RGBColor.from_string(value)


def set_font(run, size, value, bold=False):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color(value)
    props = run._r.get_or_add_rPr()
    props.set("lang", "zh-CN")
    for tag in ("a:latin", "a:ea", "a:cs"):
        node = props.find(qn(tag))
        if node is None:
            node = OxmlElement(tag)
            props.append(node)
        node.set("typeface", FONT)


def background(slide, value):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color(value)


def add_text(slide, text, x, y, w, h, size=18, value=None, bold=False,
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, margin=0.04,
             rotation=0):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.rotation = rotation
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.word_wrap = True
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_before = Pt(0)
    paragraph.space_after = Pt(0)
    paragraph.line_spacing = 1.05
    run = paragraph.add_run()
    run.text = text
    set_font(run, size, value or C["ink"], bold)
    return box


def add_bullets(slide, items, x, y, w, h, size=15, value=None, accent=None,
                line_spacing=1.16):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(0.02)
    frame.margin_right = Inches(0.02)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(8)
        paragraph.line_spacing = line_spacing
        run = paragraph.add_run()
        run.text = "● "
        set_font(run, size - 2, accent or C["coral"], True)
        run = paragraph.add_run()
        run.text = item
        set_font(run, size, value or C["ink"], False)
    return box


def add_shape(slide, kind, x, y, w, h, fill, line=None, radius=True):
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = color(fill)
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = color(line)
        shape.line.width = Pt(1)
    return shape


def add_line(slide, x, y, w, h=0, value=None, width=1):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(max(h, 0.012)))
    line.fill.solid()
    line.fill.fore_color.rgb = color(value or C["border"])
    line.line.fill.background()
    return line


def add_title(slide, title, subtitle=None, dark=False):
    title_color = C["light"] if dark else C["ink"]
    add_text(slide, title, 0.72, 0.36, 11.7, 0.58, 28, title_color, True, valign=MSO_ANCHOR.MIDDLE)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0.72, 1.03, 0.72, 0.055, C["coral"])
    if subtitle:
        add_text(slide, subtitle, 1.62, 0.91, 10.7, 0.32, 12.5, C["muted"], False, valign=MSO_ANCHOR.MIDDLE)


def add_footer(slide, page, dark=False):
    value = "8A9E90" if dark else C["muted"]
    add_text(slide, "A5 景区导览服务 AI 数字人  ·  灵山小向导·灵曦", 0.72, 7.12, 8.7, 0.2, 9, value)
    add_text(slide, f"{page:02d}", 12.15, 7.08, 0.45, 0.25, 10, C["coral"], True, PP_ALIGN.RIGHT)


def add_card(slide, x, y, w, h, title, body, accent=None, dark=False,
             title_size=17, body_size=12.5):
    fill = C["dark"] if dark else C["paper"]
    line = accent or (C["divider"] if dark else C["border"])
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill, line)
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, 0.06, h, accent or C["coral"])
    add_text(slide, title, x + 0.22, y + 0.18, w - 0.4, 0.38, title_size,
             C["light"] if dark else C["ink"], True)
    add_text(slide, body, x + 0.22, y + 0.66, w - 0.42, h - 0.82, body_size,
             "D5DFD8" if dark else C["muted"], False)


def add_metric(slide, x, y, w, h, number, label, accent=None, dark=False, suffix=None):
    fill = C["section"] if dark else C["paper"]
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill,
              C["divider"] if dark else C["border"])
    add_text(slide, number, x + 0.15, y + 0.16, w - 0.3, 0.72, 31,
             accent or C["coral"], True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    if suffix:
        add_text(slide, suffix, x + 0.15, y + 0.81, w - 0.3, 0.24, 10,
                 C["muted"], False, PP_ALIGN.CENTER)
    add_text(slide, label, x + 0.15, y + h - 0.49, w - 0.3, 0.28, 11.5,
             C["light"] if dark else C["muted"], True, PP_ALIGN.CENTER)


def add_picture_fit(slide, path, x, y, w, h):
    with Image.open(path) as image:
        iw, ih = image.size
    ratio = min(w / iw, h / ih)
    width = iw * ratio
    height = ih * ratio
    left = x + (w - width) / 2
    top = y + (h - height) / 2
    return slide.shapes.add_picture(str(path), Inches(left), Inches(top), Inches(width), Inches(height))


def add_tag(slide, text, x, y, w, fill=None, value=None):
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, 0.34, fill or C["divider"])
    add_text(slide, text, x + 0.05, y + 0.03, w - 0.1, 0.25, 10.5,
             value or C["light"], True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)


def new_slide(prs, fill=C["cream"]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background(slide, fill)
    return slide


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
prs.core_properties.title = "灵山小向导·灵曦——景区导览服务 AI 数字人"
prs.core_properties.subject = "第十五届中国软件杯 A5 产品方案介绍"
prs.core_properties.author = "[请填写团队名称]"
prs.core_properties.keywords = "AI数字人,景区导览,RAG,LiveTalking,GLM"


slide = new_slide(prs, C["dark"])
add_shape(slide, MSO_SHAPE.OVAL, 9.55, -1.35, 5.7, 5.7, None, C["divider"])
add_shape(slide, MSO_SHAPE.OVAL, 10.05, -0.85, 4.7, 4.7, None, C["divider"])
add_shape(slide, MSO_SHAPE.RECTANGLE, 0.82, 0.72, 0.06, 0.62, C["coral"])
add_text(slide, "第十五届中国软件杯 · A5 赛题", 1.03, 0.74, 5.5, 0.34, 13, "B7C8BE", True)
add_text(slide, "灵山小向导", 0.82, 1.66, 7.1, 0.92, 49, C["light"], True)
add_text(slide, "让灵曦成为每位游客的专属导游", 0.82, 2.62, 7.2, 0.72, 29, C["coral"], True)
add_text(slide, "景区导览服务 AI 数字人", 0.85, 3.52, 5.8, 0.38, 17, "C9D7CF", False)
add_tag(slide, "多模态交互", 0.85, 4.28, 1.55, C["divider"])
add_tag(slide, "可靠 RAG", 2.55, 4.28, 1.35, C["divider"])
add_tag(slide, "数字人讲解", 4.05, 4.28, 1.55, C["divider"])
add_tag(slide, "运营闭环", 5.75, 4.28, 1.35, C["divider"])
add_text(slide, "参赛学校：[待填写]   团队：[待填写]   指导教师：[待填写]", 0.85, 6.53, 7.3, 0.32, 12, "9CB0A4")
add_shape(slide, MSO_SHAPE.OVAL, 9.58, 0.62, 2.95, 6.05, C["section"], C["divider"])
add_picture_fit(slide, AVATAR, 9.58, 0.72, 2.95, 5.9)
add_text(slide, "灵曦", 10.42, 6.45, 1.35, 0.32, 15, C["light"], True, PP_ALIGN.CENTER)


slide = new_slide(prs)
add_title(slide, "四个行业痛点，必须同时解决", "功能完整度 40 分 · 技术创新 30 分 · 体验 20 分 · 文档 10 分")
pain_cards = [
    ("01", "导游资源稀缺", "旺季供不应求\n服务时间受限", "7×24 小时数字人", C["coral"]),
    ("02", "信息单向传递", "录音内容固定\n无法回答个性问题", "语音/文字/图片互动", C["sage"]),
    ("03", "缺乏情感连接", "设备冰冷\n讲解缺少亲和力", "自然语音与口型表情", C["gold"]),
    ("04", "管理反馈盲区", "关注点不可量化\n服务难持续优化", "实时数据与游客洞察", C["teal"]),
]
for index, (num, title, pain, answer, accent) in enumerate(pain_cards):
    x = 0.72 + index * 3.12
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, 1.55, 2.78, 4.72, C["paper"], C["border"])
    add_text(slide, num, x + 0.2, 1.75, 0.65, 0.52, 24, accent, True)
    add_text(slide, title, x + 0.2, 2.32, 2.28, 0.45, 17, C["ink"], True)
    add_line(slide, x + 0.2, 2.94, 2.3, value=C["border"])
    add_text(slide, pain, x + 0.2, 3.18, 2.3, 1.06, 14, C["muted"], False)
    add_shape(slide, MSO_SHAPE.DOWN_ARROW, x + 1.11, 4.34, 0.56, 0.52, accent)
    add_text(slide, answer, x + 0.2, 5.05, 2.38, 0.85, 14.5, accent, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
add_text(slide, "产品不是一个聊天框，而是一条从游客服务到景区决策的数据闭环。", 1.0, 6.53, 11.3, 0.38, 17, C["ink"], True, PP_ALIGN.CENTER)
add_footer(slide, 2)


slide = new_slide(prs)
add_title(slide, "一次游览，形成一条持续优化的体验闭环", "游客端负责服务，管理端把每次互动转化为运营依据")
journey = [
    ("01", "选择偏好", "历史 / 自然 / 亲子"),
    ("02", "语音提问", "文字与麦克风统一入口"),
    ("03", "数字人讲解", "引用、语音、口型、表情"),
    ("04", "多模态导览", "识景、路线、弱定位"),
    ("05", "游客反馈", "景点评分与文字建议"),
    ("06", "运营优化", "知识更新与服务建议"),
]
for index, (num, title, body) in enumerate(journey):
    x = 0.62 + index * 2.08
    accent = [C["coral"], C["sage"], C["gold"], C["teal"], C["green"], C["coral"]][index]
    add_shape(slide, MSO_SHAPE.OVAL, x + 0.63, 1.64, 0.66, 0.66, accent)
    add_text(slide, num, x + 0.66, 1.8, 0.6, 0.24, 11, C["white"], True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    if index < len(journey) - 1:
        add_line(slide, x + 1.31, 1.96, 1.42, 0.03, C["border"])
    add_text(slide, title, x, 2.54, 1.92, 0.42, 15.5, C["ink"], True, PP_ALIGN.CENTER)
    add_text(slide, body, x, 3.03, 1.92, 0.86, 11.8, C["muted"], False, PP_ALIGN.CENTER)
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.82, 4.38, 11.68, 1.7, C["dark"], C["divider"])
add_text(slide, "游客实时交互", 1.08, 4.67, 2.25, 0.4, 18, C["light"], True)
add_text(slide, "问答 · 位置 · 反馈 · 情绪", 1.08, 5.17, 2.45, 0.32, 12, "B6C8BE")
add_shape(slide, MSO_SHAPE.CHEVRON, 3.62, 4.83, 0.72, 0.56, C["coral"])
add_text(slide, "运营洞察", 4.64, 4.67, 2.05, 0.4, 18, C["light"], True)
add_text(slide, "热点 · 趋势 · 建议 · 客群", 4.64, 5.17, 2.45, 0.32, 12, "B6C8BE")
add_shape(slide, MSO_SHAPE.CHEVRON, 7.05, 4.83, 0.72, 0.56, C["gold"])
add_text(slide, "内容与服务优化", 8.05, 4.67, 2.55, 0.4, 18, C["light"], True)
add_text(slide, "知识库 · 路线 · 现场引导", 8.05, 5.17, 2.65, 0.32, 12, "B6C8BE")
add_shape(slide, MSO_SHAPE.CIRCULAR_ARROW, 11.0, 4.65, 0.85, 0.85, C["teal"])
add_footer(slide, 3)


slide = new_slide(prs, C["dark"])
add_title(slide, "五层架构，把实时体验与运营分析解耦", "内部服务最小暴露 · 云端模型可替换 · GPU 按需使用", True)
layers = [
    ("交互层", "游客网页 / Android App / 管理后台", C["coral"]),
    ("网关层", "FastAPI · HTTPS · SSE · WebRTC · 鉴权", C["gold"]),
    ("AI 编排层", "GLM-ASR · GLM-4V · GLM-TTS · RAG/GLM", C["sage"]),
    ("数字人层", "LiveTalking · Wav2Lip FP16 · Live2D 回退", C["teal"]),
    ("数据洞察层", "SQLite · 14 万历史记录 · HumanOmni", C["green"]),
]
for index, (name, body, accent) in enumerate(layers):
    y = 1.47 + index * 1.0
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1.08, y, 11.15, 0.76, C["section"], C["divider"])
    add_shape(slide, MSO_SHAPE.RECTANGLE, 1.08, y, 1.82, 0.76, accent)
    add_text(slide, name, 1.18, y + 0.17, 1.6, 0.3, 15, C["white"], True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    add_text(slide, body, 3.2, y + 0.15, 8.62, 0.34, 14.5, C["light"], False, valign=MSO_ANCHOR.MIDDLE)
add_tag(slide, "8001/8443 游客端", 1.08, 6.62, 2.0, C["divider"])
add_tag(slide, "8010 LiveTalking", 3.28, 6.62, 1.85, C["divider"])
add_tag(slide, "8020 RAG（仅本机）", 5.33, 6.62, 2.08, C["divider"])
add_tag(slide, "8444 独立管理端", 7.61, 6.62, 1.92, C["divider"])
add_tag(slide, "密钥不下发", 9.73, 6.62, 1.55, C["divider"])
add_footer(slide, 4, True)


slide = new_slide(prs)
add_title(slide, "可靠 RAG：回答有依据，资料不足就拒答", "BGE-M3 + FAISS + 景点元数据过滤 + 有界会话历史")
steps = [
    ("文档切片", "官方 DOCX / TXT / XLSX\n管理员上传资料"),
    ("向量编码", "BGE-M3\n1024 维归一化向量"),
    ("候选过滤", "景点名 / 别名\n景区元数据缩小范围"),
    ("FAISS 检索", "IndexFlatIP\n余弦相似度 Top-K"),
    ("受约束生成", "事实必须来自上下文\n返回稳定引用或明确拒答"),
]
for index, (title, body) in enumerate(steps):
    x = 0.66 + index * 2.06
    add_card(slide, x, 1.56, 1.82, 2.28, f"{index + 1:02d}  {title}", body,
             [C["coral"], C["sage"], C["gold"], C["teal"], C["green"]][index],
             title_size=13.5, body_size=10.8)
    if index < 4:
        add_shape(slide, MSO_SHAPE.CHEVRON, x + 1.84, 2.39, 0.34, 0.55, C["border"])
add_metric(slide, 0.82, 4.42, 2.8, 1.74, "15/15", "标准事实题", C["coral"])
add_metric(slide, 3.88, 4.42, 2.8, 1.74, "100%", "事实问答准确率", C["green"])
add_metric(slide, 6.94, 4.42, 2.8, 1.74, "97", "当前知识片段", C["teal"])
add_metric(slide, 10.0, 4.42, 2.5, 1.74, "180s", "RAG GPU 空闲释放", C["gold"])
add_text(slide, "CPU 常驻协调器 → 检索时选择空闲 GPU → 连续问答复用 → 空闲后 worker 退出并彻底释放 CUDA", 1.02, 6.44, 11.25, 0.36, 13.5, C["ink"], True, PP_ALIGN.CENTER)
add_footer(slide, 5)


slide = new_slide(prs, C["section"])
add_title(slide, "首句即说：跨模型流水线压缩用户等待", "不等待全文完成，按完整语义句进入 TTS 与口型驱动", True)
add_picture_fit(slide, AVATAR, 0.66, 1.44, 2.75, 5.45)
pipeline = [
    ("01", "GLM-ASR", "语音转文字", C["coral"]),
    ("02", "RAG / GLM", "流式生成", C["gold"]),
    ("03", "语义分句", "完整短句", C["sage"]),
    ("04", "GLM-TTS", "PCM 合并", C["teal"]),
    ("05", "Wav2Lip", "口型同步", C["green"]),
]
for index, (num, title, body, accent) in enumerate(pipeline):
    x = 3.42 + index * 1.75
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, 2.0, 1.48, 1.56, C["dark"], C["divider"])
    add_text(slide, num, x + 0.12, 2.16, 0.38, 0.25, 10.5, accent, True)
    add_text(slide, title, x + 0.12, 2.5, 1.24, 0.35, 14, C["light"], True, PP_ALIGN.CENTER)
    add_text(slide, body, x + 0.12, 3.02, 1.24, 0.25, 10.5, "A8BBB0", False, PP_ALIGN.CENTER)
    if index < 4:
        add_shape(slide, MSO_SHAPE.CHEVRON, x + 1.49, 2.5, 0.26, 0.5, accent)
add_metric(slide, 3.44, 4.33, 2.46, 1.64, "3373", "语音结束→首个数字人音频", C["coral"], True, "ms · 单次实测")
add_metric(slide, 6.14, 4.33, 2.46, 1.64, "3618", "质量优先热链路首响", C["gold"], True, "ms · 三轮中位数")
add_metric(slide, 8.84, 4.33, 2.46, 1.64, "1.0×", "自然原速播报", C["teal"], True, "完整语义段")
add_text(slide, "再次提问会取消旧播报队列；LiveTalking 不可用时回退 Live2D，文字问答不中断。", 3.46, 6.35, 7.95, 0.4, 13.2, "C4D2CA", True, PP_ALIGN.CENTER)
add_footer(slide, 6, True)


slide = new_slide(prs)
add_title(slide, "多模态与个性化，让导览从“能回答”走向“懂场景”", "所有模态最终回到景区知识与游客真实意图")
multi = [
    ("拍照识景", "图片质量检测 + CLIP 景点图集召回 + GLM-4V 复核反证；RAG 对确认后景点提供有依据讲解，不校准识别结果。", C["coral"]),
    ("分众路线", "历史、自然、亲子三类偏好，结合时间与位置输出顺路路线和讲解重点。", C["sage"]),
    ("弱定位降级", "GPS 不可用时切换景点码、Wi-Fi/二维码适配层或手动选择。", C["gold"]),
    ("游客感受度", "HumanOmni 七类情绪 + 文本方面 + 显式评分，保留模式、置信度和降级状态。", C["teal"]),
]
for index, (title, body, accent) in enumerate(multi):
    x = 0.78 + (index % 2) * 6.14
    y = 1.5 + (index // 2) * 2.43
    add_card(slide, x, y, 5.65, 2.05, title, body, accent, title_size=18, body_size=13)
    add_shape(slide, MSO_SHAPE.OVAL, x + 4.85, y + 0.28, 0.46, 0.46, accent)
    add_text(slide, ["图", "路", "位", "感"][index], x + 4.94, y + 0.37, 0.28, 0.22, 11, C["white"], True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1.52, 6.34, 10.28, 0.5, C["dark"])
add_text(slide, "原则：不把历史满意度伪装成情绪标签；模型不可用时明确显示降级，不阻塞实时问答。", 1.68, 6.43, 9.98, 0.26, 12.5, C["light"], True, PP_ALIGN.CENTER)
add_footer(slide, 7)


slide = new_slide(prs)
add_title(slide, "游客端：一个页面完成问答、路线、识景与定位", "真实运行界面 · HTTPS 麦克风 · WebRTC 数字人 · 轻量模式自动回退")
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.55, 1.34, 9.12, 5.46, C["soft"], C["border"])
add_picture_fit(slide, ASSETS / "visitor.png", 0.68, 1.47, 8.86, 5.2)
add_card(slide, 9.92, 1.36, 2.72, 1.12, "一问即答", "文字与麦克风统一进入同一条问答链路", C["coral"], title_size=15, body_size=10.5)
add_card(slide, 9.92, 2.68, 2.72, 1.12, "四类工具", "问答 / 路线 / 识景 / 定位", C["sage"], title_size=15, body_size=10.5)
add_card(slide, 9.92, 4.0, 2.72, 1.12, "自然讲解", "流式语音、口型与情绪状态同步", C["gold"], title_size=15, body_size=10.5)
add_card(slide, 9.92, 5.32, 2.72, 1.12, "可用性优先", "数字人暂不可用时仍可先提问", C["teal"], title_size=15, body_size=10.5)
add_footer(slide, 8)


slide = new_slide(prs)
add_title(slide, "管理后台：把真实互动转化为景区运营动作", "实时数据、知识维护、数字人配置、游客感受度与历史洞察")
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.55, 1.34, 9.12, 5.46, C["soft"], C["border"])
add_picture_fit(slide, ASSETS / "admin.png", 0.68, 1.47, 8.86, 5.2)
add_metric(slide, 9.92, 1.4, 2.72, 1.18, "176", "今日服务游客", C["coral"])
add_metric(slide, 9.92, 2.76, 2.72, 1.18, "235", "累计问答", C["teal"])
add_metric(slide, 9.92, 4.12, 2.72, 1.18, "3182", "平均响应 ms", C["gold"])
add_metric(slide, 9.92, 5.48, 2.72, 1.18, "140,447", "历史游客记录", C["green"])
add_footer(slide, 9)


slide = new_slide(prs, C["dark"])
add_title(slide, "用可复现数据证明：准确、快速、稳定、可解释", "所有数字注明时间、环境和测试口径，不以演示效果代替验证", True)
metrics = [
    ("15/15", "事实冒烟基线", "关键词预检；≥80 题冻结集赛前补齐", C["coral"]),
    ("3373ms", "语音首响", "低于赛题 5 秒", C["gold"]),
    ("35", "API 回归", "全部 passed", C["sage"]),
    ("97", "知识片段", "3 类来源", C["teal"]),
]
for index, (number, label, sub, accent) in enumerate(metrics):
    add_metric(slide, 0.78 + index * 3.12, 1.52, 2.76, 1.78, number, label, accent, True, sub)
rows = [
    ("事实问答", "≥ 90%", "100%", "通过"),
    ("语音问答首响", "< 5 秒", "3.373 秒", "通过"),
    ("多模态核心模型", "至少 1 个", "GLM-4V / HumanOmni", "通过"),
    ("本地知识库", "必须构建", "BGE-M3 + FAISS", "通过"),
]
headers = ["评测项", "赛题门槛", "当前结果", "结论"]
widths = [3.05, 2.25, 4.2, 1.25]
x0, y0 = 1.28, 4.02
cursor = x0
for text_value, width in zip(headers, widths):
    add_shape(slide, MSO_SHAPE.RECTANGLE, cursor, y0, width, 0.5, C["coral"])
    add_text(slide, text_value, cursor + 0.04, y0 + 0.1, width - 0.08, 0.26, 12, C["white"], True, PP_ALIGN.CENTER)
    cursor += width
for row_index, row in enumerate(rows):
    cursor = x0
    y = y0 + 0.5 + row_index * 0.55
    for col_index, (text_value, width) in enumerate(zip(row, widths)):
        fill = C["section"] if row_index % 2 == 0 else C["dark"]
        add_shape(slide, MSO_SHAPE.RECTANGLE, cursor, y, width, 0.55, fill, C["divider"])
        add_text(slide, text_value, cursor + 0.04, y + 0.12, width - 0.08, 0.26, 11.2,
                 C["green"] if col_index == 3 else C["light"], col_index == 3, PP_ALIGN.CENTER)
        cursor += width
add_footer(slide, 10, True)


slide = new_slide(prs)
add_title(slide, "七个创新点，最终落到可复制的景区价值", "技术不是堆模型，而是围绕准确性、自然度、可用性与运营闭环做取舍")
innovations = [
    ("首句即说", "跨模型按语义句流水线", C["coral"]),
    ("视觉校准", "识景结果二次进入 RAG", C["sage"]),
    ("事实追溯", "引用、过滤、拒答共同约束", C["gold"]),
    ("按需 GPU", "闲时释放，连续问答复用", C["teal"]),
    ("弱定位", "GPS→景点码→手动多级降级", C["green"]),
    ("双层洞察", "实时互动 + 14 万历史记录", C["coral"]),
    ("证据化情绪", "模态、置信度、降级可解释", C["sage"]),
]
for index, (title, body, accent) in enumerate(innovations):
    col = index % 4
    row = index // 4
    x = 0.66 + col * 3.08
    y = 1.5 + row * 1.62
    width = 2.78 if index < 4 else 3.72
    if row == 1:
        x = 0.94 + col * 4.08
    add_card(slide, x, y, width, 1.32, title, body, accent, title_size=15, body_size=10.8)
add_text(slide, "规模化路线", 0.82, 5.05, 1.55, 0.36, 16, C["ink"], True)
roadmap = [
    ("01", "景区试点", "完善 POI、可信证书与现场定位", C["coral"]),
    ("02", "多景区复制", "多租户、方言、IP 头像与内容运营", C["gold"]),
    ("03", "边云协同", "量化模型、边缘缓存与旺季容量治理", C["teal"]),
]
for index, (num, title, body, accent) in enumerate(roadmap):
    x = 2.42 + index * 3.48
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, 4.76, 3.12, 1.35, C["dark"], C["divider"])
    add_text(slide, num, x + 0.16, 4.95, 0.45, 0.28, 11, accent, True)
    add_text(slide, title, x + 0.66, 4.92, 2.12, 0.32, 14.5, C["light"], True)
    add_text(slide, body, x + 0.18, 5.44, 2.74, 0.38, 10.8, "B7C8BE", False, PP_ALIGN.CENTER)
add_footer(slide, 11)


slide = new_slide(prs, C["dark"])
add_shape(slide, MSO_SHAPE.OVAL, 9.56, -1.1, 5.2, 5.2, None, C["divider"])
add_text(slide, "让每位游客拥有专属导游", 0.82, 0.8, 8.8, 0.72, 34, C["light"], True)
add_text(slide, "让每次互动成为运营依据", 0.82, 1.58, 8.8, 0.72, 34, C["coral"], True)
add_line(slide, 0.85, 2.55, 5.45, 0.05, C["divider"])
add_text(slide, "团队信息（提交前补充）", 0.85, 2.9, 4.2, 0.4, 18, C["light"], True)
team = [
    "队长 / 产品架构：[姓名]",
    "大模型与 RAG：[姓名]",
    "数字人与语音：[姓名]",
    "前后端与数据：[姓名]",
    "指导教师：[姓名]",
]
add_bullets(slide, team, 0.86, 3.48, 4.9, 2.45, 14, "C9D7CF", C["coral"], 1.05)
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 6.35, 3.0, 5.75, 2.95, C["section"], C["divider"])
add_text(slide, "答辩现场建议", 6.68, 3.28, 2.3, 0.4, 18, C["light"], True)
add_bullets(slide, [
    "现场语音提问并展示引用与口型",
    "演示路线、识景、定位与反馈",
    "刷新后台证明真实数据闭环",
    "最后展示 15/15 与 3373ms 测试证据",
], 6.68, 3.85, 4.85, 1.7, 13, "C9D7CF", C["gold"], 1.0)
add_tag(slide, "谢谢 · Q&A", 9.65, 6.44, 1.82, C["coral"], C["white"])
add_footer(slide, 12, True)


prs.save(OUT)
print(OUT)
